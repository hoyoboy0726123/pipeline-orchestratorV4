"""
人工確認節點的「上一步驟輸出檔案預覽」：
把檔案 render 成 PNG 傳給 Telegram，讓使用者在手機上直接看到內容不用 SSH 回電腦。

設計原則：
  B1（主要）= pure Python + headless 函式庫，不開 Excel/Word 等 GUI App
            輸出「資料結構」忠實度 60-80%，樣式會掉但資料都在
  B2（後備）= libreoffice --headless 轉 PDF → pypdfium2 render
            輸出「版式」忠實度 80-90%，要求使用者裝 LibreOffice

呼叫：
  render_file_preview(file_path, out_dir) -> list[str]
    → 回傳一至多個 PNG 路徑（多頁 PDF / 多張投影片會回多張）
    → 失敗時回空 list，caller 自行處理（例如傳個「預覽失敗」的錯誤截圖）
"""
from __future__ import annotations

import logging
import shutil
import subprocess
import time
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)


# ── 單頁圖片尺寸（統一規格，避免 TG 上大小不一致）──
PREVIEW_WIDTH = 1600   # 合適手機閱讀
PREVIEW_HEIGHT = 1200
# 多頁文件最多 render 幾頁（PDF / PPTX），超過截斷
MAX_PAGES = 3
# 表格類型（xlsx/csv）顯示前幾列
MAX_TABLE_ROWS = 20
# 文字類型（docx/txt）顯示前幾行
MAX_TEXT_LINES = 40


def render_file_preview(file_path: str, out_dir: Optional[str] = None) -> list[str]:
    """主要入口。回傳 PNG 路徑 list（失敗回 [])。"""
    p = Path(file_path)
    if not p.exists():
        log.warning(f"[preview] 檔案不存在：{file_path}")
        return []
    if not p.is_file():
        log.warning(f"[preview] 不是檔案：{file_path}")
        return []

    out = Path(out_dir) if out_dir else p.parent
    out.mkdir(parents=True, exist_ok=True)
    ext = p.suffix.lower().lstrip(".")

    # B1 路線：每種格式有對應 renderer
    try:
        if ext in ("png", "jpg", "jpeg", "gif", "webp", "bmp"):
            return _render_image(p, out)
        if ext == "pdf":
            return _render_pdf(p, out)
        if ext in ("xlsx", "xls", "xlsm"):
            return _render_spreadsheet(p, out)
        if ext == "csv":
            return _render_csv(p, out)
        if ext == "docx":
            return _render_docx(p, out)
        if ext == "pptx":
            return _render_pptx(p, out)
        if ext in ("txt", "md", "log", "json", "yaml", "yml", "html", "xml"):
            return _render_text(p, out)
    except Exception as e:
        log.warning(f"[preview] B1 render 失敗（{ext}）：{e}；改試 LibreOffice 後備")

    # B2 後備：LibreOffice → PDF → pypdfium2
    try:
        return _render_via_libreoffice(p, out)
    except Exception as e:
        log.warning(f"[preview] LibreOffice 後備也失敗：{e}")

    # 都失敗 → 回一張「不支援」的提示圖
    return _render_unsupported(p, out)


# ── B1 路線：各格式 renderer ───────────────────────────────────────────

def _render_image(src: Path, out: Path) -> list[str]:
    """圖片：直接複製（或轉成 PNG 並縮到合理尺寸）。"""
    from PIL import Image
    im = Image.open(src)
    if im.mode not in ("RGB", "L"):
        im = im.convert("RGB")
    # 縮到預覽尺寸內（保持長寬比）
    im.thumbnail((PREVIEW_WIDTH, PREVIEW_HEIGHT), Image.LANCZOS)
    dst = out / f"{src.stem}_preview.png"
    im.save(dst, "PNG", optimize=True)
    return [str(dst)]


def _render_pdf(src: Path, out: Path) -> list[str]:
    """PDF：pypdfium2 直接 render page，結果跟 PDF 原樣（因為是 PDF engine）。
    要顯式 close 不然 Windows 會 hold file handle 一陣子、後續刪檔會失敗。
    """
    import pypdfium2 as pdfium
    pdf = pdfium.PdfDocument(src)
    try:
        n_pages = min(len(pdf), MAX_PAGES)
        results: list[str] = []
        for i in range(n_pages):
            page = pdf[i]
            try:
                img = page.render(scale=2).to_pil()
                dst = out / f"{src.stem}_preview_page{i+1}.png"
                img.save(dst, "PNG", optimize=True)
                results.append(str(dst))
            finally:
                page.close()
        return results
    finally:
        pdf.close()


def _render_spreadsheet(src: Path, out: Path) -> list[str]:
    """Excel：逐 sheet 讀前幾列、各產一張 table PNG。
    局限：不保留 cell 顏色 / 合併 / 格式化 / 圖表（這些要走 LibreOffice 後備）。
    """
    import pandas as pd
    # sheet_name=None → 回 dict {sheet_name: DataFrame}，讀所有 sheet
    sheets_dict = pd.read_excel(src, sheet_name=None, nrows=MAX_TABLE_ROWS)
    if not sheets_dict:
        # 退回單張空圖以免回空 list
        return [_render_text_to_png(["（Excel 無 sheet）"], src, out, title=src.name)]
    results: list[str] = []
    sheet_names = list(sheets_dict.keys())
    # 超過 MAX_PAGES 張 sheet 只取前 N 張，最後補一張「還有 X 張」提示
    show_sheets = sheet_names[:MAX_PAGES]
    for idx, sn in enumerate(show_sheets, start=1):
        df = sheets_dict[sn]
        # 檔名加 sheet 序號避免檔名衝突；title 顯示 sheet 名方便辨識
        title = f"{src.name} — Sheet: {sn}（{len(sheets_dict)} 張中的第 {idx} 張，前 {len(df)} 列）"
        png = _render_dataframe(df, src, out, title=title, suffix=f"_sheet{idx}")
        results.append(png)
    if len(sheet_names) > MAX_PAGES:
        more = sheet_names[MAX_PAGES:]
        results.append(_render_text_to_png(
            [f"... 還有 {len(more)} 張 sheet 未顯示（MAX_PAGES={MAX_PAGES}）:", "",
             *[f"  • {n}" for n in more[:20]]],
            src, out, title=f"{src.name} — 更多 sheets", suffix="_more",
        ))
    return results


def _render_csv(src: Path, out: Path) -> list[str]:
    """CSV：同 spreadsheet。試 utf-8 / utf-8-sig / big5 三種編碼。"""
    import pandas as pd
    df = None
    last_err: Optional[Exception] = None
    for enc in ("utf-8-sig", "utf-8", "big5"):
        try:
            df = pd.read_csv(src, nrows=MAX_TABLE_ROWS, encoding=enc)
            break
        except UnicodeDecodeError as e:
            last_err = e
            continue
    if df is None:
        raise RuntimeError(f"無法用 utf-8/utf-8-sig/big5 解碼 CSV：{last_err}")
    return [_render_dataframe(df, src, out, title=f"{src.name}（前 {len(df)} 列）")]


def _render_docx(src: Path, out: Path) -> list[str]:
    """docx：python-docx 抽前幾段文字，用 PIL 畫。
    局限：保留不了粗體/顏色/表格/圖片；要真實版式走 LibreOffice 後備。
    """
    from docx import Document
    doc = Document(src)
    lines: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # 依段落長度略做換行（手機閱讀）
            while len(text) > 60:
                lines.append(text[:60])
                text = text[60:]
            lines.append(text)
        if len(lines) >= MAX_TEXT_LINES:
            break
    if not lines:
        lines = ["（文件無文字內容 — 可能全是圖片或表格，建議用 LibreOffice 預覽）"]
    elif len(lines) >= MAX_TEXT_LINES:
        lines.append(f"... (已截斷，只顯示前 {MAX_TEXT_LINES} 行)")
    return [_render_text_to_png(lines, src, out, title=f"{src.name}（docx 文字摘要）")]


def _render_pptx(src: Path, out: Path) -> list[str]:
    """pptx：python-pptx 抽每張投影片的文字內容。
    局限：看不到真實投影片版式、圖片、圖表；要看版式用 LibreOffice 後備。
    """
    from pptx import Presentation
    prs = Presentation(src)
    results: list[str] = []
    for i, slide in enumerate(prs.slides):
        if i >= MAX_PAGES:
            break
        lines = [f"=== 投影片 {i+1}/{len(prs.slides)} ==="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        lines.append(t)
        if len(lines) == 1:
            lines.append("（本張無文字；只有圖片或圖表）")
        png = _render_text_to_png(lines, src, out,
                                  title=f"{src.name} — 第 {i+1} 張",
                                  suffix=f"_slide{i+1}")
        results.append(png)
    if len(prs.slides) > MAX_PAGES:
        results.append(_render_text_to_png(
            [f"... 還有 {len(prs.slides) - MAX_PAGES} 張投影片未顯示（MAX_PAGES={MAX_PAGES}）"],
            src, out, title=f"{src.name} — 更多", suffix="_more",
        ))
    return results


def _render_text(src: Path, out: Path) -> list[str]:
    """純文字類（txt/md/json/yaml/html...）：直接讀前幾行畫成 PNG。"""
    try:
        content = src.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        content = f"[讀取失敗：{e}]"
    lines = content.splitlines()[:MAX_TEXT_LINES]
    if not lines:
        lines = ["（空檔案）"]
    elif len(content.splitlines()) > MAX_TEXT_LINES:
        lines.append(f"... (已截斷，完整 {len(content.splitlines())} 行)")
    return [_render_text_to_png(lines, src, out, title=f"{src.name}（前 {len(lines)} 行）")]


# ── Render 輔助（共用 PIL 畫圖） ─────────────────────────────────────

def _render_dataframe(df, src: Path, out: Path, title: str, suffix: str = "") -> str:
    """pandas DataFrame → matplotlib table PNG。
    matplotlib 的 ax.table 相對陽春，但保證不閃屏、不需 Office。
    suffix：多 sheet 時用來避免檔名衝突（例：_sheet1 / _sheet2）
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    # 欄太多會擠爆；保留前 8 欄，其餘標示 "…"
    cols = list(df.columns)
    show_cols = cols[:8]
    truncated_cols = len(cols) > 8
    display_df = df[show_cols].copy()
    if truncated_cols:
        display_df["…"] = "…"
    # 全部轉字串（數字 / NaN 都用 str 避免渲染問題）
    display_df = display_df.astype(str)
    # matplotlib 的中文字型：選常見幾個試
    for font in ["Microsoft JhengHei", "Microsoft YaHei", "PingFang HK", "SimHei", "Arial Unicode MS"]:
        try:
            matplotlib.font_manager.findfont(font, fallback_to_default=False)
            matplotlib.rcParams["font.family"] = font
            break
        except Exception:
            pass
    fig_h = max(3, 0.45 * (len(display_df) + 2))
    fig_w = min(24, max(10, 1.6 * len(display_df.columns)))
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.axis("off")
    ax.set_title(title, fontsize=14, pad=10)
    tbl = ax.table(
        cellText=display_df.values,
        colLabels=display_df.columns,
        loc="center",
        cellLoc="left",
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)
    tbl.scale(1, 1.4)
    # 標題列粉淺灰底
    for j in range(len(display_df.columns)):
        tbl[(0, j)].set_facecolor("#e5e7eb")
        tbl[(0, j)].set_text_props(weight="bold")
    dst = out / f"{src.stem}_preview{suffix}.png"
    fig.savefig(dst, bbox_inches="tight", dpi=130, facecolor="white")
    plt.close(fig)
    return str(dst)


def _render_text_to_png(lines: list[str], src: Path, out: Path, title: str,
                        suffix: str = "") -> str:
    """幾行文字 → PNG（PIL）。純 PIL 避開 matplotlib 文字層級 bug。"""
    from PIL import Image, ImageDraw, ImageFont
    # 選中文字型（Windows 內建微軟正黑體；找不到就 fallback default）
    font_size = 18
    title_size = 22
    line_h = int(font_size * 1.55)
    font = _load_cjk_font(font_size)
    font_title = _load_cjk_font(title_size, bold=True)

    # 估算需要的高度
    body_h = len(lines) * line_h + 40
    height = max(400, body_h + 80)
    width = PREVIEW_WIDTH
    im = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(im)

    # 標題
    draw.rectangle([(0, 0), (width, 50)], fill="#1f2937")
    draw.text((20, 12), title, fill="white", font=font_title)

    # 內容
    y = 70
    for ln in lines:
        draw.text((25, y), ln, fill="#111827", font=font)
        y += line_h
    dst = out / f"{src.stem}_preview{suffix}.png"
    im.save(dst, "PNG", optimize=True)
    return str(dst)


def _load_cjk_font(size: int, bold: bool = False):
    """找一個能顯示中文的 TrueType font。Windows 優先微軟正黑體；
    找不到就 PIL 預設 bitmap font（中文會變方塊，但至少不 crash）。"""
    from PIL import ImageFont
    candidates = [
        "C:\\Windows\\Fonts\\msjh.ttc" if not bold else "C:\\Windows\\Fonts\\msjhbd.ttc",
        "C:\\Windows\\Fonts\\msyh.ttc",
        "C:\\Windows\\Fonts\\simhei.ttf",
        "/System/Library/Fonts/PingFang.ttc",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def _render_unsupported(src: Path, out: Path) -> list[str]:
    """所有方法都失敗 → 回一張錯誤提示圖，讓使用者看到「預覽失敗」而不是一片黑。"""
    lines = [
        "⚠️ 無法產生檔案預覽",
        "",
        f"檔案：{src.name}",
        f"大小：{src.stat().st_size:,} bytes",
        f"副檔名：{src.suffix}",
        "",
        "可能原因：",
        "  • 格式不支援（非 xlsx/csv/docx/pptx/pdf/圖片/純文字）",
        "  • 檔案損壞或加密",
        "  • 本機未裝 LibreOffice（後備機制）",
        "",
        "請到後端主機直接查看此檔案。",
    ]
    return [_render_text_to_png(lines, src, out, title="預覽失敗", suffix="_unsupported")]


# ── B2 後備：LibreOffice 無頭轉 PDF ───────────────────────────────────

def _libreoffice_binary() -> Optional[str]:
    """找 LibreOffice 可執行檔。Windows 裝了通常在 Program Files。
    找不到回 None，caller 放棄 B2。"""
    # PATH 上的 soffice
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    # Windows 常見安裝路徑
    import os
    for root in (os.environ.get("ProgramFiles", ""), os.environ.get("ProgramFiles(x86)", "")):
        if not root:
            continue
        guess = Path(root) / "LibreOffice" / "program" / "soffice.exe"
        if guess.exists():
            return str(guess)
    return None


def _render_via_libreoffice(src: Path, out: Path) -> list[str]:
    """用 LibreOffice 無頭模式把檔案轉 PDF，再用 pypdfium2 render。
    比 B1 貼近真實版式（~80-90% 還原），但要使用者裝 LibreOffice。"""
    binary = _libreoffice_binary()
    if not binary:
        raise RuntimeError("找不到 LibreOffice（soffice）— 請到 libreoffice.org 下載安裝")
    # 轉 PDF 到臨時目錄，避免弄髒 out
    import tempfile
    with tempfile.TemporaryDirectory() as tmp:
        log.info(f"[preview] LibreOffice 轉檔：{src.name} → PDF ({binary})")
        t0 = time.time()
        result = subprocess.run(
            [binary, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(src)],
            capture_output=True, text=True, timeout=60,
        )
        dur = time.time() - t0
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice 轉檔失敗 rc={result.returncode}: {result.stderr[:300]}")
        pdf_path = Path(tmp) / f"{src.stem}.pdf"
        if not pdf_path.exists():
            raise RuntimeError(f"LibreOffice 轉完了但 PDF 不存在：{pdf_path}")
        log.info(f"[preview] LibreOffice 轉檔完成（{dur:.1f}s），render PDF...")
        # 搬 PDF 到 out 再 render（讓 out 下有穩定路徑，debug 時好看）
        stable_pdf = out / f"{src.stem}_libre.pdf"
        shutil.copyfile(pdf_path, stable_pdf)
    return _render_pdf(stable_pdf, out)
